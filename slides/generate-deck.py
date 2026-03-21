#!/usr/bin/env python3
"""Generate Dubai Dip tokenization structure slide deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
BG_DARK = RGBColor(0x0a, 0x0a, 0x0f)
BG_CARD = RGBColor(0x14, 0x14, 0x1e)
BG_CARD2 = RGBColor(0x1a, 0x1a, 0x28)
ACCENT_GREEN = RGBColor(0x10, 0xb9, 0x81)
ACCENT_PURPLE = RGBColor(0x8b, 0x5c, 0xf6)
ACCENT_BLUE = RGBColor(0x38, 0xbd, 0xf8)
ACCENT_AMBER = RGBColor(0xf5, 0x9e, 0x0b)
ACCENT_RED = RGBColor(0xef, 0x44, 0x44)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY = RGBColor(0x94, 0xa3, 0xb8)
LIGHT_GRAY = RGBColor(0xcb, 0xd5, 0xe1)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Smaller corner radius
    shape.adjustments[0] = 0.03
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txbox

def add_multiline(slide, left, top, width, height, lines, font_size=16, color=WHITE, line_spacing=1.3):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, (text, sz, clr, bld) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.space_after = Pt(4)
    return txbox

def add_arrow(slide, left, top, width, height, color=ACCENT_GREEN):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_right_arrow(slide, left, top, width, height, color=ACCENT_GREEN):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

add_text(slide, Inches(1.5), Inches(2), Inches(13), Inches(1.5),
         "DUBAI DIP", 72, ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(3.5), Inches(13), Inches(1),
         "Tokenized UAE Equities on Solana", 36, WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(4.5), Inches(13), Inches(0.8),
         "Structure & Distribution Options", 24, GRAY, align=PP_ALIGN.CENTER)

# Bottom bar
add_rect(slide, Inches(0), Inches(8.2), Inches(16), Inches(0.03), ACCENT_GREEN)
add_text(slide, Inches(1.5), Inches(8.3), Inches(5), Inches(0.5),
         "dubaidip.xyz", 14, GRAY)
add_text(slide, Inches(10.5), Inches(8.3), Inches(4), Inches(0.5),
         "Confidential — March 2026", 14, GRAY, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 2: THE OPPORTUNITY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(1), Inches(0.5), Inches(14), Inches(0.8),
         "THE OPPORTUNITY", 14, ACCENT_GREEN, bold=True)
add_text(slide, Inches(1), Inches(1), Inches(14), Inches(1),
         "UAE equities are at crisis-driven discounts.", 40, WHITE, bold=True)
add_text(slide, Inches(1), Inches(1.8), Inches(14), Inches(0.8),
         "Global crypto-native investors want exposure but can't access DFM/ADX easily.", 22, GRAY)

# Stats row
stats = [
    ("16", "UAE Equities\nTracked"),
    ("-15% to -36%", "Drawdown from\nPre-Crisis ATH"),
    ("$0", "Minimum for\nOnchain Access"),
    ("24/7", "Trading on\nSolana"),
]
for i, (val, label) in enumerate(stats):
    x = Inches(1 + i * 3.6)
    card = add_rect(slide, x, Inches(3), Inches(3.2), Inches(2.2), BG_CARD, RGBColor(0x2a, 0x2a, 0x3a))
    add_text(slide, x + Inches(0.3), Inches(3.3), Inches(2.6), Inches(1),
             val, 36, ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.3), Inches(4.2), Inches(2.6), Inches(0.8),
             label, 16, GRAY, align=PP_ALIGN.CENTER)

# Problem statement
add_rect(slide, Inches(1), Inches(5.8), Inches(14), Inches(2.2), BG_CARD)
add_multiline(slide, Inches(1.4), Inches(6), Inches(13.2), Inches(2), [
    ("The Problem", 18, ACCENT_AMBER, True),
    ("Opening a UAE brokerage account requires residency, visa, Emirates ID. International brokers (IBKR, Saxo)", 16, LIGHT_GRAY, False),
    ("have limited DFM/ADX coverage and high minimums. Crypto investors sitting on USDC have zero access.", 16, LIGHT_GRAY, False),
    ("", 8, GRAY, False),
    ("The Solution: Tokenize the underlying equities on Solana. Buy with USDC. Instant settlement. No broker needed.", 16, WHITE, True),
])

# ============================================================
# SLIDE 3: PRECEDENT — ONDO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(1), Inches(0.5), Inches(14), Inches(0.8),
         "PRECEDENT", 14, ACCENT_GREEN, bold=True)
add_text(slide, Inches(1), Inches(1), Inches(14), Inches(1),
         "Ondo Finance proved the model.", 40, WHITE, bold=True)

# Ondo structure
add_rect(slide, Inches(1), Inches(2.2), Inches(6.8), Inches(5.5), BG_CARD, RGBColor(0x2a, 0x2a, 0x3a))
add_multiline(slide, Inches(1.4), Inches(2.4), Inches(6), Inches(5), [
    ("Ondo Global Markets Structure", 18, ACCENT_BLUE, True),
    ("", 6, GRAY, False),
    ("Issuer:  BVI SPV", 16, WHITE, False),
    ("Governing Law:  Swiss", 16, WHITE, False),
    ("Custody:  Regulated broker-dealer", 16, WHITE, False),
    ("Chains:  Ethereum → BNB → Solana", 16, WHITE, False),
    ("Assets:  200+ US stocks & ETFs", 16, WHITE, False),
    ("", 6, GRAY, False),
    ("ADGM Approval (Mar 3, 2026)", 16, ACCENT_GREEN, True),
    ("First-ever tokenized securities approval in Abu Dhabi.", 14, GRAY, False),
    ("Allows UAE institutions to trade via Binance's MTF.", 14, GRAY, False),
    ("", 6, GRAY, False),
    ("Key insight: ADGM is the distribution approval layer.", 14, ACCENT_AMBER, True),
    ("The issuer entity sits offshore (BVI). UAE is just where", 14, GRAY, False),
    ("institutions access it — not where it's created.", 14, GRAY, False),
])

# Our opportunity
add_rect(slide, Inches(8.2), Inches(2.2), Inches(6.8), Inches(5.5), BG_CARD, RGBColor(0x2a, 0x2a, 0x3a))
add_multiline(slide, Inches(8.6), Inches(2.4), Inches(6), Inches(5), [
    ("Dubai Dip — Same Playbook", 18, ACCENT_GREEN, True),
    ("", 6, GRAY, False),
    ("Issuer:  BVI SPV (to be formed)", 16, WHITE, False),
    ("Governing Law:  BVI / English", 16, WHITE, False),
    ("Custody:  Regulated UAE broker", 16, WHITE, False),
    ("Chain:  Solana", 16, WHITE, False),
    ("Assets:  16 UAE equities (DFM + ADX)", 16, WHITE, False),
    ("", 6, GRAY, False),
    ("What's Different", 16, ACCENT_AMBER, True),
    ("→ Underlying = UAE equities (not US)", 14, WHITE, False),
    ("→ Thesis-driven: crisis discount", 14, WHITE, False),
    ("→ Crypto-native distribution first", 14, WHITE, False),
    ("→ No UAE exchange needed for v1", 14, WHITE, False),
    ("→ Institutional (ADGM) wrapper = Phase 2", 14, WHITE, False),
])

# ============================================================
# SLIDE 4: OPTION A — OFFSHORE-FIRST (RECOMMENDED)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(1), Inches(0.5), Inches(14), Inches(0.8),
         "OPTION A — RECOMMENDED", 14, ACCENT_GREEN, bold=True)
add_text(slide, Inches(1), Inches(1), Inches(14), Inches(1),
         "Offshore-First: BVI SPV + Solana Distribution", 36, WHITE, bold=True)
add_text(slide, Inches(1), Inches(1.7), Inches(14), Inches(0.6),
         "Fastest path. No UAE licensing required. Ondo-proven structure.", 20, GRAY)

# Flow diagram - 5 boxes with arrows
boxes = [
    ("1. FORM SPV", "BVI Company\n\nIssues tokenized\nnotes/certificates\nrepresenting UAE\nequity exposure", ACCENT_PURPLE),
    ("2. RAISE", "Investor USDC\n→ SPV\n\nEligible non-UAE\ninvestors subscribe\nvia dubaidip.xyz", ACCENT_BLUE),
    ("3. BUY", "SPV → Broker\n→ DFM/ADX\n\nSPV purchases\nunderlying equities\nthrough regulated\nbroker", ACCENT_GREEN),
    ("4. TOKENIZE", "Mint on Solana\n\nSPL tokens issued\n1:1 against held\nequities. Oracle\nfeeds live NAV.", ACCENT_AMBER),
    ("5. TRADE", "Global Access\n\nHolders trade tokens\n24/7 on Solana DEXs.\nRedeem for underlying\nvia SPV.", ACCENT_GREEN),
]

for i, (title, desc, color) in enumerate(boxes):
    x = Inches(0.5 + i * 3.1)
    card = add_rect(slide, x, Inches(2.6), Inches(2.7), Inches(4.2), BG_CARD, color)
    add_text(slide, x + Inches(0.2), Inches(2.8), Inches(2.3), Inches(0.6),
             title, 15, color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), Inches(3.5), Inches(2.3), Inches(3), desc, 13, LIGHT_GRAY, align=PP_ALIGN.CENTER)
    # Arrow between boxes
    if i < 4:
        add_right_arrow(slide, x + Inches(2.75), Inches(4.4), Inches(0.3), Inches(0.3), RGBColor(0x3a, 0x3a, 0x4a))

# Pros/cons
add_rect(slide, Inches(0.5), Inches(7.2), Inches(7.2), Inches(1.3), BG_CARD)
add_multiline(slide, Inches(0.8), Inches(7.3), Inches(6.6), Inches(1.2), [
    ("✅  Advantages", 14, ACCENT_GREEN, True),
    ("No UAE license needed • Fastest launch • Global reach • Low regulatory overhead • Proven by Ondo", 13, LIGHT_GRAY, False),
])

add_rect(slide, Inches(8.3), Inches(7.2), Inches(7.2), Inches(1.3), BG_CARD)
add_multiline(slide, Inches(8.6), Inches(7.3), Inches(6.6), Inches(1.2), [
    ("⚠️  Considerations", 14, ACCENT_AMBER, True),
    ("BVI setup ~$5-10K • Need custodial broker with DFM/ADX access • UAE residents may be geo-fenced", 13, LIGHT_GRAY, False),
])

# ============================================================
# SLIDE 5: OPTION B — UAE EXCHANGE PARTNER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(1), Inches(0.5), Inches(14), Inches(0.8),
         "OPTION B — UAE EXCHANGE PARTNER", 14, ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(1), Inches(14), Inches(1),
         "Partner with a Regulated UAE Exchange", 36, WHITE, bold=True)
add_text(slide, Inches(1), Inches(1.7), Inches(14), Inches(0.6),
         "Higher credibility. Institutional access. Slower path.", 20, GRAY)

# Three sub-options
partners = [
    ("Bybit", "Federal UAE Approval", [
        "Bybit holds federal-level UAE approval",
        "Could list tokenized UAE equities on their platform",
        "Large existing crypto user base",
        "Would need commercial partnership deal",
        "Bybit handles compliance / KYC",
    ], ACCENT_BLUE),
    ("Bahrain Exchange", "CBB Licensed", [
        "Friendly exchange already holds",
        "Central Bank of Bahrain approval",
        "Existing crypto-securities framework",
        "Faster than getting new UAE license",
        "Bahrain = lighter regulatory touch",
    ], ACCENT_PURPLE),
    ("Regulated Broker Clients", "Existing Licensees", [
        "Lawyer's clients: licensed brokers/exchanges",
        "Already have distribution infrastructure",
        "Could white-label the product",
        "Brings existing institutional client base",
        "Revenue share model possible",
    ], ACCENT_GREEN),
]

for i, (name, subtitle, bullets, color) in enumerate(partners):
    x = Inches(0.8 + i * 5)
    card = add_rect(slide, x, Inches(2.5), Inches(4.6), Inches(4.5), BG_CARD, color)
    add_text(slide, x + Inches(0.3), Inches(2.7), Inches(4), Inches(0.6),
             name, 24, color, bold=True)
    add_text(slide, x + Inches(0.3), Inches(3.2), Inches(4), Inches(0.4),
             subtitle, 14, GRAY)
    for j, bullet in enumerate(bullets):
        add_text(slide, x + Inches(0.3), Inches(3.8 + j * 0.5), Inches(4), Inches(0.5),
                 f"→  {bullet}", 13, LIGHT_GRAY)

# Bottom note
add_rect(slide, Inches(0.8), Inches(7.4), Inches(14.4), Inches(1), BG_CARD)
add_multiline(slide, Inches(1.1), Inches(7.5), Inches(13.8), Inches(0.9), [
    ("💡  Option B is about distribution, not issuance.", 15, ACCENT_AMBER, True),
    ("The BVI SPV still issues the tokens. A UAE/Bahrain exchange partner provides the regulated distribution channel for institutional investors.", 13, LIGHT_GRAY, False),
])

# ============================================================
# SLIDE 6: OPTION C — HYBRID
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(1), Inches(0.5), Inches(14), Inches(0.8),
         "OPTION C — HYBRID (RECOMMENDED PATH)", 14, ACCENT_GREEN, bold=True)
add_text(slide, Inches(1), Inches(1), Inches(14), Inches(1),
         "Launch Offshore, Add UAE Distribution Later", 36, WHITE, bold=True)

# Timeline
phases = [
    ("PHASE 1: NOW", "Weeks 1-4", [
        "Form BVI SPV ($5-10K, 2 weeks)",
        "Open custodial broker account",
        "Build Solana token program",
        "Soft launch via dubaidip.xyz waitlist",
        "Seed round: raise initial USDC",
    ], ACCENT_GREEN, "Launch with crypto-native\ndistribution only"),
    ("PHASE 2: SCALE", "Months 2-4", [
        "Onboard 100+ holders",
        "Add more UAE equities",
        "Integrate Solana DEX liquidity",
        "Pyth/Switchboard oracle for NAV",
        "Approach exchange partners",
    ], ACCENT_BLUE, "Prove demand with\nonchain traction"),
    ("PHASE 3: INSTITUTIONAL", "Months 4-8", [
        "Partner with UAE/Bahrain exchange",
        "ADGM distribution approval",
        "Family office / SWF access",
        "Expand to GCC equities",
        "Revenue: mgmt fee + spread",
    ], ACCENT_PURPLE, "Add regulated wrapper\nfor institutional AUM"),
]

for i, (title, timeline, bullets, color, note) in enumerate(phases):
    x = Inches(0.5 + i * 5.2)
    card = add_rect(slide, x, Inches(2.2), Inches(4.8), Inches(5), BG_CARD, color)
    add_text(slide, x + Inches(0.3), Inches(2.4), Inches(4.2), Inches(0.5),
             title, 18, color, bold=True)
    add_text(slide, x + Inches(0.3), Inches(2.9), Inches(4.2), Inches(0.4),
             timeline, 14, GRAY)
    for j, bullet in enumerate(bullets):
        add_text(slide, x + Inches(0.3), Inches(3.5 + j * 0.45), Inches(4.2), Inches(0.4),
                 f"•  {bullet}", 13, LIGHT_GRAY)
    # Note at bottom
    add_rect(slide, x + Inches(0.2), Inches(6), Inches(4.4), Inches(0.9), BG_CARD2)
    add_text(slide, x + Inches(0.4), Inches(6.1), Inches(4), Inches(0.8),
             note, 12, color, align=PP_ALIGN.CENTER)

# Arrow flow
add_right_arrow(slide, Inches(5.35), Inches(4.5), Inches(0.35), Inches(0.3), RGBColor(0x3a, 0x3a, 0x4a))
add_right_arrow(slide, Inches(10.55), Inches(4.5), Inches(0.35), Inches(0.3), RGBColor(0x3a, 0x3a, 0x4a))

# ============================================================
# SLIDE 7: KEY QUESTIONS / NEXT STEPS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(1), Inches(0.5), Inches(14), Inches(0.8),
         "KEY DECISIONS & NEXT STEPS", 14, ACCENT_GREEN, bold=True)
add_text(slide, Inches(1), Inches(1), Inches(14), Inches(1),
         "What needs to be decided now.", 40, WHITE, bold=True)

# Left: Key decisions
add_rect(slide, Inches(0.8), Inches(2.2), Inches(7), Inches(5.8), BG_CARD)
add_multiline(slide, Inches(1.1), Inches(2.4), Inches(6.4), Inches(5.5), [
    ("🔑  Key Decisions", 20, ACCENT_AMBER, True),
    ("", 6, GRAY, False),
    ("1. Custodial Broker", 16, WHITE, True),
    ("Who holds the DFM/ADX shares for the SPV? Need a broker with", 13, GRAY, False),
    ("UAE market access that's comfortable with a crypto-linked entity.", 13, GRAY, False),
    ("", 6, GRAY, False),
    ("2. NAV & Redemption Mechanism", 16, WHITE, True),
    ("How do token holders redeem for underlying? Ondo uses authorized", 13, GRAY, False),
    ("participants (like ETFs). We need a similar mechanism.", 13, GRAY, False),
    ("", 6, GRAY, False),
    ("3. Price Oracle", 16, WHITE, True),
    ("Who feeds DFM/ADX prices onchain? Pyth? Switchboard? Custom?", 13, GRAY, False),
    ("DFM/ADX close at 3 PM — what happens with 24/7 token trading?", 13, GRAY, False),
    ("", 6, GRAY, False),
    ("4. Regulatory Counsel", 16, WHITE, True),
    ("BVI counsel for SPV formation + token legal opinion.", 13, GRAY, False),
    ("Confirm tokens classify as 'notes' not 'securities' in BVI.", 13, GRAY, False),
])

# Right: Next steps
add_rect(slide, Inches(8.2), Inches(2.2), Inches(7), Inches(5.8), BG_CARD)
add_multiline(slide, Inches(8.5), Inches(2.4), Inches(6.4), Inches(5.5), [
    ("⚡  Immediate Next Steps", 20, ACCENT_GREEN, True),
    ("", 6, GRAY, False),
    ("□  Get BVI counsel quote for SPV formation", 16, WHITE, False),
    ("    Timeline: 1-2 weeks once engaged", 13, GRAY, False),
    ("", 6, GRAY, False),
    ("□  Identify custodial broker", 16, WHITE, False),
    ("    Explore: Bahrain exchange, IBKR, Saxo", 13, GRAY, False),
    ("", 6, GRAY, False),
    ("□  Draft Solana token program spec", 16, WHITE, False),
    ("    SPL token + mint authority + oracle integration", 13, GRAY, False),
    ("", 6, GRAY, False),
    ("□  Legal opinion on token classification", 16, WHITE, False),
    ("    Required before any token issuance", 13, GRAY, False),
    ("", 6, GRAY, False),
    ("□  Seed investor outreach", 16, WHITE, False),
    ("    dubaidip.xyz waitlist → first subscribers", 13, GRAY, False),
    ("", 6, GRAY, False),
    ("□  Approach exchange partners (Phase 2)", 16, WHITE, False),
    ("    Bybit / Bahrain exchange / lawyer's clients", 13, GRAY, False),
])

# Bottom
add_rect(slide, Inches(0), Inches(8.5), Inches(16), Inches(0.03), ACCENT_GREEN)
add_text(slide, Inches(1), Inches(8.55), Inches(14), Inches(0.4),
         "dubaidip.xyz  •  Confidential  •  March 2026", 12, GRAY, align=PP_ALIGN.CENTER)

# Save
output_path = "/Users/clawdbot/.openclaw/workspace/buy-the-dubai-dip/slides/dubai-dip-tokenization.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
